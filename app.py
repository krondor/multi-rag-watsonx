import os
import streamlit as st
import tempfile
import pandas as pd
import json
import xml.etree.ElementTree as ET
import yaml
from bs4 import BeautifulSoup
from pptx import Presentation
from docx import Document

from langchain.document_loaders import PyPDFLoader, TextLoader
from langchain.indexes import VectorstoreIndexCreator
from langchain.chains import RetrievalQA
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate

from ibm_watson_machine_learning.foundation_models import Model
from ibm_watson_machine_learning.foundation_models.extensions.langchain import WatsonxLLM
from ibm_watson_machine_learning.metanames import GenTextParamsMetaNames as GenParams
from ibm_watson_machine_learning.foundation_models.utils.enums import DecodingMethods

# Initialize index and chain to None1
index = None
rag_chain = None

# Custom loader for DOCX files
class DocxLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        document = Document(self.file_path)
        text_content = [para.text for para in document.paragraphs]
        return " ".join(text_content)

# Custom loader for PPTX files
class PptxLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        presentation = Presentation(self.file_path)
        text_content = [shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")]
        return " ".join(text_content)

# Custom loader for additional file types
def load_csv(file_path):
    df = pd.read_csv(file_path)
    page_size = 100
    page_number = st.number_input("Page number", min_value=1, max_value=(len(df) // page_size) + 1, step=1, value=1)
    start_index = (page_number - 1) * page_size
    end_index = start_index + page_size
    st.dataframe(df.iloc[start_index:end_index])
    return df.to_string(index=False)

def load_json(file_path):
    with open(file_path, 'r') as file:
        data = json.load(file)
    return json.dumps(data, indent=2)

def load_xml(file_path):
    tree = ET.parse(file_path)
    root = tree.getroot()
    return ET.tostring(root, encoding="unicode")

def load_yaml(file_path):
    with open(file_path, 'r') as file:
        data = yaml.safe_load(file)
    return yaml.dump(data)

def load_html(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        soup = BeautifulSoup(file, 'html.parser')
    return soup.get_text()

# Caching function to load various file types
@st.cache_resource
def load_file(file_name, file_type):
    loaders = []
    text = None

    if file_type == "pdf":
        loaders = [PyPDFLoader(file_name)]
    elif file_type == "docx":
        loader = DocxLoader(file_name)
        text = loader.load()
    elif file_type == "pptx":
        loader = PptxLoader(file_name)
        text = loader.load()
    elif file_type == "txt":
        loaders = [TextLoader(file_name)]
    elif file_type == "csv":
        text = load_csv(file_name)
    elif file_type == "json":
        text = load_json(file_name)
    elif file_type == "xml":
        text = load_xml(file_name)
    elif file_type == "yaml":
        text = load_yaml(file_name)
    elif file_type == "html":
        text = load_html(file_name)
    elif file_type == "htm":
        text = load_html(file_name)    
    else:
        st.error("Unsupported file type.")
        return None

    if text:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".txt") as temp_file:
            temp_file.write(text.encode("utf-8"))
            temp_file_path = temp_file.name
        loaders = [TextLoader(temp_file_path)]

    if loaders:
        index = VectorstoreIndexCreator(
            embedding=HuggingFaceEmbeddings(model_name="all-MiniLM-L12-v2"),
            text_splitter=RecursiveCharacterTextSplitter(chunk_size=450, chunk_overlap=50)
        ).from_loaders(loaders)
        st.success("Index created successfully!")
        return index
    return None

# Watsonx API setup
watsonx_api_key =  os.getenv("WATSONX_API_KEY")
watsonx_project_id = os.getenv("WATSONX_PROJECT_ID")

if not watsonx_api_key or not watsonx_project_id:
    st.error("API Key or Project ID is not set. Please set them as environment variables.")

prompt_template_br = PromptTemplate(
    input_variables=["context", "question"],
    template="""<|begin_of_text|><|start_header_id|>system<|end_header_id|>
I am a helpful assistant.
<|eot_id|>
{context}
<|start_header_id|>user<|end_header_id|>
{question}<|eot_id|>
"""
)

with st.sidebar:
    st.title("Multi-Document Retrieval with Watsonx")
    st.sidebar.write("") 
    st.sidebar.markdown("Developed by **Abdul Rahman H**")
    watsonx_model = st.selectbox("Model", ["meta-llama/llama-3-405b-instruct", "codellama/codellama-34b-instruct-hf", "ibm/granite-20b-multilingual"])
    max_new_tokens = st.slider("Max output tokens", min_value=100, max_value=4000, value=600, step=100)
    decoding_method = st.radio("Decoding", (DecodingMethods.GREEDY.value, DecodingMethods.SAMPLE.value))
    parameters = {
        GenParams.DECODING_METHOD: decoding_method,
        GenParams.MAX_NEW_TOKENS: max_new_tokens,
        GenParams.MIN_NEW_TOKENS: 1,
        GenParams.TEMPERATURE: 0,
        GenParams.TOP_K: 50,
        GenParams.TOP_P: 1,
        GenParams.STOP_SEQUENCES: [],
        GenParams.REPETITION_PENALTY: 1
    }
    st.info("Upload a file to use RAG")
    uploaded_file = st.file_uploader("Upload file", accept_multiple_files=False, type=["pdf", "docx", "txt", "pptx", "csv", "json", "xml", "yaml", "html"])

    if uploaded_file is not None:
        bytes_data = uploaded_file.read()
        st.write("Filename:", uploaded_file.name)
        with open(uploaded_file.name, 'wb') as f:
            f.write(bytes_data)
        file_type = uploaded_file.name.split('.')[-1].lower()
        index = load_file(uploaded_file.name, file_type)

    model_name = watsonx_model

st.info("Setting up Watsonx...")
my_credentials = {
    "url": "https://us-south.ml.cloud.ibm.com",
    "apikey": watsonx_api_key
}
params = parameters
project_id = watsonx_project_id
space_id = None
verify = False
model = WatsonxLLM(model=Model(model_name, my_credentials, params, project_id, space_id, verify))

if model:
    st.info(f"Model {model_name} ready.")
    chain = LLMChain(llm=model, prompt=prompt_template_br, verbose=True)

if chain and index is not None:
    rag_chain = RetrievalQA.from_chain_type(
        llm=model,
        chain_type="stuff",
        retriever=index.vectorstore.as_retriever(),
        chain_type_kwargs={"prompt": prompt_template_br},
        return_source_documents=False,
        verbose=True
    )
    st.info("Document-based retrieval is ready.")
else:
    st.warning("No document uploaded or chain setup issue.")

# Chat loop
if "messages" not in st.session_state:
    st.session_state.messages = []

for message in st.session_state.messages:
    st.chat_message(message["role"]).markdown(message["content"])

prompt = st.chat_input("Ask your question here", disabled=False if chain else True)

if prompt:
    st.chat_message("user").markdown(prompt)
    if rag_chain:
        response_text = rag_chain.run(prompt).strip()
    else:
        response_text = chain.run(question=prompt, context="").strip()
        
    st.session_state.messages.append({'role': 'User', 'content': prompt})
    st.chat_message("assistant").markdown(response_text)
    st.session_state.messages.append({'role': 'Assistant', 'content': response_text})

